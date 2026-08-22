#!/usr/bin/env python3
"""Build day-mode PowerPoint template for Ahmed Alharbi brand."""

from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "pptx" / "قالب-نهاري.pptx"
ASSETS = ROOT / "png" / "شعارات"

# Brand tokens — day mode
C_BG = RGBColor(0xF2, 0xF4, 0xF6)
C_TEXT = RGBColor(0x0B, 0x0D, 0x10)
C_TEAL = RGBColor(0x4F, 0xD8, 0xC4)
C_GRAY = RGBColor(0x8A, 0x93, 0xA0)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BORDER = RGBColor(0xE2, 0xE5, 0xE9)

FONT = "Thmanyah Sans"
FONT_MONO = "Courier New"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
M = Inches(0.65)
FOOTER_Y = Inches(6.85)
CONTENT_W = SLIDE_W - M * 2


def rtl(paragraph) -> None:
    p_pr = paragraph._element.get_or_add_pPr()
    p_pr.set("rtl", "1")


def style_run(run, *, size: int, bold: bool = False, color=C_TEXT, name: str = FONT) -> None:
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_para(
    text_frame,
    text: str,
    *,
    size: int = 24,
    bold: bool = False,
    color=C_TEXT,
    align=PP_ALIGN.RIGHT,
    name: str = FONT,
    space_after: int = 6,
) -> None:
    p = text_frame.paragraphs[0] if not text_frame.text else text_frame.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    rtl(p)
    run = p.add_run()
    run.text = text
    style_run(run, size=size, bold=bold, color=color, name=name)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    lines: list[tuple[str, dict]],
    *,
    valign=MSO_ANCHOR.TOP,
) -> object:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.clear()
    first = True
    for text, opts in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = opts.get("align", PP_ALIGN.RIGHT)
        p.space_after = Pt(opts.get("space_after", 8))
        rtl(p)
        run = p.add_run()
        run.text = text
        style_run(
            run,
            size=opts.get("size", 24),
            bold=opts.get("bold", False),
            color=opts.get("color", C_TEXT),
            name=opts.get("name", FONT),
        )
    return box


def add_round_rect(slide, left, top, width, height, *, fill=C_WHITE, line=C_BORDER, line_w=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_w)
    shape.adjustments[0] = 0.12
    return shape


def add_teal_bar(slide, left, top, *, width=Inches(0.55), height=Inches(0.07)):
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_TEAL
    bar.line.fill.background()
    bar.adjustments[0] = 0.5
    return bar


def add_pill(slide, left, top, text: str, *, filled=True):
    w, h = Inches(2.2), Inches(0.42)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    pill.fill.solid()
    pill.line.width = Pt(1.5)
    if filled:
        pill.fill.fore_color.rgb = C_TEAL
        pill.line.fill.background()
        txt_color = C_TEXT
    else:
        pill.fill.fore_color.rgb = RGBColor(0xEC, 0xFA, 0xF7)
        pill.line.color.rgb = C_TEAL
        txt_color = C_TEAL
    pill.adjustments[0] = 0.5
    tf = pill.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    rtl(p)
    run = p.add_run()
    run.text = text
    style_run(run, size=14, bold=True, color=txt_color)
    return pill


def add_takeaway_card(slide, left, top, width, text: str):
    card = add_round_rect(slide, left, top, width, Inches(0.95), fill=C_WHITE, line=C_BORDER)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left + width - Inches(0.08), top + Inches(0.12), Inches(0.06), Inches(0.71)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = C_TEAL
    accent.line.fill.background()
    add_textbox(
        slide,
        left + Inches(0.25),
        top + Inches(0.18),
        width - Inches(0.45),
        Inches(0.6),
        [(text, {"size": 20, "align": PP_ALIGN.CENTER})],
        valign=MSO_ANCHOR.MIDDLE,
    )
    return card


def add_header(slide, *, badge: str | None = None):
    mark = ASSETS / "logo-mark-on-light-1024.png"
    if mark.exists():
        slide.shapes.add_picture(str(mark), SLIDE_W - M - Inches(0.55), M - Inches(0.05), height=Inches(0.55))

    header = slide.shapes.add_textbox(M, M - Inches(0.05), Inches(4.5), Inches(0.55))
    tf = header.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    rtl(p)
    run = p.add_run()
    run.text = "احمد الحربي"
    style_run(run, size=18, bold=True, color=C_TEXT)

    if badge:
        add_pill(slide, SLIDE_W - M - Inches(2.35), M + Inches(0.05), badge, filled=False)


def add_footer(slide, page: str):
    add_textbox(
        slide,
        M,
        FOOTER_Y,
        Inches(1.2),
        Inches(0.35),
        [(page, {"size": 13, "color": C_GRAY, "name": FONT, "align": PP_ALIGN.LEFT})],
        valign=MSO_ANCHOR.MIDDLE,
    )

    handle_w = Inches(2.4)
    pill = add_round_rect(
        slide,
        (SLIDE_W - handle_w) / 2,
        FOOTER_Y - Inches(0.02),
        handle_w,
        Inches(0.38),
        fill=C_WHITE,
        line=C_BORDER,
    )
    tf = pill.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "@Ahmed014x"
    style_run(run, size=14, bold=False, color=C_TEAL, name=FONT_MONO)

    sig_x = SLIDE_W - M - Inches(0.95)
    bar_a = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sig_x, FOOTER_Y + Inches(0.04), Inches(0.72), Inches(0.1))
    bar_a.fill.solid()
    bar_a.fill.fore_color.rgb = C_TEAL
    bar_a.line.fill.background()
    bar_a.adjustments[0] = 0.4
    bar_b = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, sig_x + Inches(0.28), FOOTER_Y + Inches(0.18), Inches(0.72), Inches(0.1)
    )
    bar_b.fill.solid()
    bar_b.fill.fore_color.rgb = RGBColor(0xD5, 0xD9, 0xDE)
    bar_b.line.fill.background()
    bar_b.adjustments[0] = 0.4


def add_bg_picture(slide, bg_path: Path) -> None:
    pic = slide.shapes.add_picture(str(bg_path), 0, 0, width=SLIDE_W, height=SLIDE_H)
    sp_tree = slide.shapes._spTree
    element = pic._element
    sp_tree.remove(element)
    sp_tree.insert(2, element)


def set_slide_bg(slide, bg_path: Path) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C_BG
    add_bg_picture(slide, bg_path)


def make_bg_png(path: Path) -> None:
    w, h = 1920, 1080
    img = Image.new("RGB", (w, h), "#F2F4F6")
    draw = ImageDraw.Draw(img)
    # subtle dot grid
    for y in range(0, h, 44):
        for x in range(0, w, 44):
            draw.ellipse((x, y, x + 2, y + 2), fill=(210, 214, 218))
    # teal glow top-right
    for i in range(120, 0, -4):
        alpha = int(18 * (i / 120))
        color = (242 - alpha // 3, 250 - alpha // 6, 246)
        draw.ellipse((w - 900 - i * 3, -200 - i, w + 200, 500 + i), fill=color)
    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path, "PNG")


def build() -> None:
    bg_path = ROOT / "pptx" / "_assets" / "slide-bg-day.png"
    make_bg_png(bg_path)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ── 1. Cover ──
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, bg_path)
    add_header(slide)
    add_pill(slide, (SLIDE_W - Inches(2.4)) / 2, Inches(2.0), "عنوان فرعي", filled=True)
    add_textbox(
        slide,
        M,
        Inches(2.65),
        CONTENT_W,
        Inches(1.4),
        [("عنوان العرض", {"size": 54, "bold": True, "align": PP_ALIGN.CENTER})],
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        M,
        Inches(4.05),
        CONTENT_W,
        Inches(0.6),
        [("احمد الحربي  ·  مهندس ذكاء اصطناعي", {"size": 22, "color": C_GRAY, "align": PP_ALIGN.CENTER})],
    )
    add_footer(slide, "1")

    # ── 2. Section ──
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, bg_path)
    add_header(slide, badge="قسم")
    add_textbox(
        slide,
        SLIDE_W - M - Inches(1.2),
        Inches(2.35),
        Inches(1.0),
        Inches(0.9),
        [("01", {"size": 56, "bold": True, "color": C_TEAL})],
    )
    add_teal_bar(slide, M, Inches(3.35), width=Inches(0.7))
    add_textbox(
        slide,
        M,
        Inches(3.55),
        CONTENT_W,
        Inches(1.0),
        [("عنوان القسم", {"size": 44, "bold": True})],
    )
    add_textbox(
        slide,
        M,
        Inches(4.55),
        Inches(6.5),
        Inches(0.6),
        [("سطر يوضّح وش بنغطّي في هذا الجزء.", {"size": 20, "color": C_GRAY})],
    )
    add_footer(slide, "2")

    # ── 3. Content + bullets ──
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, bg_path)
    add_header(slide, badge="محتوى")
    add_textbox(slide, M, Inches(1.35), Inches(4), Inches(0.35), [("هوك قصير", {"size": 16, "bold": True, "color": C_TEAL})])
    add_textbox(slide, M, Inches(1.75), CONTENT_W, Inches(0.9), [("عنوان الشريحة", {"size": 36, "bold": True})])
    add_teal_bar(slide, M, Inches(2.65))
    add_textbox(
        slide,
        M,
        Inches(2.85),
        Inches(7.5),
        Inches(0.55),
        [("سطرين رماديين تحت العنوان — مو أكثر.", {"size": 18, "color": C_GRAY})],
    )
    card = add_round_rect(slide, M, Inches(3.55), Inches(5.8), Inches(2.35))
    bullets = [
        "نقطة أولى — واضحة ومباشرة",
        "نقطة ثانية — مثال أو رقم",
        "نقطة ثالثة — الخلاصة العملية",
    ]
    add_textbox(
        slide,
        M + Inches(0.25),
        Inches(3.8),
        Inches(5.3),
        Inches(1.9),
        [(b, {"size": 20, "space_after": 14}) for b in bullets],
    )
    add_takeaway_card(slide, M, Inches(6.05), CONTENT_W, "الخلاصة = شريط واحد بحافة تركواز")
    add_footer(slide, "3")

    # ── 4. Two columns ──
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, bg_path)
    add_header(slide, badge="عمودان")
    add_textbox(slide, M, Inches(1.35), CONTENT_W, Inches(0.8), [("مقارنة أو تقسيم", {"size": 36, "bold": True})])
    add_teal_bar(slide, M, Inches(2.15))
    col_w = (CONTENT_W - Inches(0.25)) / 2
    for idx, (title, sub, accent) in enumerate(
        [
            ("الخيار الأول", "سطر رمادي يشرح الفكرة", False),
            ("الخيار الثاني", "سطر رمادي يشرح الفكرة", True),
        ]
    ):
        left = M + idx * (col_w + Inches(0.25))
        card = add_round_rect(
            slide,
            left,
            Inches(2.55),
            col_w,
            Inches(3.2),
            fill=RGBColor(0xEC, 0xFA, 0xF7) if accent else C_WHITE,
            line=C_TEAL if accent else C_BORDER,
            line_w=2 if accent else 1,
        )
        add_textbox(
            slide,
            left + Inches(0.2),
            Inches(2.85),
            col_w - Inches(0.4),
            Inches(0.5),
            [(title, {"size": 24, "bold": True, "color": C_TEAL if accent else C_TEXT, "align": PP_ALIGN.CENTER})],
        )
        add_textbox(
            slide,
            left + Inches(0.2),
            Inches(3.55),
            col_w - Inches(0.4),
            Inches(0.8),
            [(sub, {"size": 17, "color": C_GRAY, "align": PP_ALIGN.CENTER})],
        )
    add_footer(slide, "4")

    # ── 5. Cards row (3-up) ──
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, bg_path)
    add_header(slide, badge="بطاقات")
    add_textbox(slide, M, Inches(1.35), CONTENT_W, Inches(0.8), [("ثلاث نقاط متوازية", {"size": 36, "bold": True})])
    add_teal_bar(slide, M, Inches(2.15))
    card_w = (CONTENT_W - Inches(0.4)) / 3
    for idx, (title, sub) in enumerate(
        [
            ("01", "كل بطاقة = عنوان + سطر"),
            ("02", "التركواز للبارز فقط"),
            ("03", "SVG ورسوم — مو صناديق نص"),
        ]
    ):
        left = M + idx * (card_w + Inches(0.2))
        add_round_rect(slide, left, Inches(2.55), card_w, Inches(2.5))
        add_textbox(
            slide,
            left + Inches(0.15),
            Inches(2.85),
            card_w - Inches(0.3),
            Inches(0.45),
            [(title, {"size": 22, "bold": True, "color": C_TEAL})],
        )
        add_textbox(
            slide,
            left + Inches(0.15),
            Inches(3.45),
            card_w - Inches(0.3),
            Inches(0.9),
            [(sub, {"size": 17, "color": C_GRAY})],
        )
    add_takeaway_card(slide, M, Inches(5.35), CONTENT_W, "كل عنصر يشرح نفسه بصريًا + سطر ثانٍ")
    add_footer(slide, "5")

    # ── 6. Visual / diagram placeholder ──
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, bg_path)
    add_header(slide, badge="رسم")
    add_textbox(slide, M, Inches(1.35), CONTENT_W, Inches(0.8), [("مخطط أو SVG", {"size": 36, "bold": True})])
    add_teal_bar(slide, M, Inches(2.15))
    frame = add_round_rect(slide, M, Inches(2.45), CONTENT_W, Inches(3.55), fill=C_WHITE)
    add_textbox(
        slide,
        M,
        Inches(3.85),
        CONTENT_W,
        Inches(0.8),
        [("ضع رسمك أو لقطة شاشة هنا", {"size": 22, "color": C_GRAY, "align": PP_ALIGN.CENTER})],
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_takeaway_card(slide, M, Inches(6.05), CONTENT_W, "الرسم يشرح — العنوان ما يكرّر كل التفاصيل")
    add_footer(slide, "6")

    # ── 7. Closing ──
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, bg_path)
    add_header(slide)
    add_textbox(
        slide,
        M,
        Inches(2.2),
        CONTENT_W,
        Inches(1.0),
        [("شكراً — خلّنا نكمّل", {"size": 48, "bold": True, "align": PP_ALIGN.CENTER})],
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_teal_bar(slide, (SLIDE_W - Inches(0.7)) / 2, Inches(3.35), width=Inches(0.7))
    add_textbox(
        slide,
        M,
        Inches(3.65),
        CONTENT_W,
        Inches(0.5),
        [("سطر ختامي قصير.", {"size": 20, "color": C_GRAY, "align": PP_ALIGN.CENTER})],
    )
    add_pill(slide, (SLIDE_W - Inches(3.0)) / 2, Inches(4.45), "تابعني @Ahmed014x", filled=True)
    add_textbox(
        slide,
        M,
        Inches(5.35),
        CONTENT_W,
        Inches(0.4),
        [("ahmednalharbii@gmail.com", {"size": 16, "color": C_GRAY, "align": PP_ALIGN.CENTER, "name": FONT_MONO})],
    )
    add_footer(slide, "7")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
